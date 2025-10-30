import json
import argparse
from pathlib import Path
from jinja2 import Environment, FileSystemLoader

# PDF変換
try:
    import pdfkit
except ImportError:
    pdfkit = None

# PPTX変換
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

# ====== 設定 ======
TEMPLATE_DIR = Path("templates")
ASSETS_DIR = Path("assets")
OUTPUT_HTML = Path("output.html")
SLIDES_JSON = Path("slides.json")

def generate_html(slides):
    """HTMLをJinja2で生成"""
    env = Environment(loader=FileSystemLoader(TEMPLATE_DIR))
    template = env.get_template("base.html")
    return template.render(slides=slides)

def export_pdf(html_path, output_pdf):
    """HTMLをPDFに変換"""
    if not pdfkit:
        print("⚠️ pdfkitがインストールされていません。`pip install pdfkit` を実行してください。")
        return
    pdfkit.from_file(str(html_path), str(output_pdf))
    print(f"✅ PDF生成完了: {output_pdf.resolve()}")

def export_pptx(slides, output_pptx):
    """JSONスライドからPPTX生成"""
    if not Presentation:
        print("⚠️ python-pptxがインストールされていません。`pip install python-pptx` を実行してください。")
        return

    prs = Presentation()

    for slide in slides:
        layout = prs.slide_layouts[1]  # タイトル＋本文
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide["title"]
        s.placeholders[1].text = slide["body"]

    prs.save(output_pptx)
    print(f"✅ PPTX生成完了: {output_pptx.resolve()}")

def main():
    # ====== CLI引数 ======
    parser = argparse.ArgumentParser(description="AIスライド生成ツール")
    parser.add_argument("--format", choices=["html", "pdf", "pptx"], default="html",
                        help="出力形式を選択（html / pdf / pptx）")
    args = parser.parse_args()

    # ====== JSON読み込み ======
    if not SLIDES_JSON.exists():
        raise FileNotFoundError(f"❌ スライドデータが見つかりません: {SLIDES_JSON}")
    with open(SLIDES_JSON, "r", encoding="utf-8") as f:
        slides = json.load(f).get("slides", [])

    if not slides:
        raise ValueError("❌ JSONにスライドデータがありません。")

    # ====== HTML生成 ======
    rendered_html = generate_html(slides)
    OUTPUT_HTML.write_text(rendered_html, encoding="utf-8")
    print(f"✅ HTML生成完了: {OUTPUT_HTML.resolve()}")

    # ====== 出力形式に応じて分岐 ======
    if args.format == "pdf":
        export_pdf(OUTPUT_HTML, Path("output.pdf"))
    elif args.format == "pptx":
        export_pptx(slides, Path("output.pptx"))
    else:
        print("💡 ブラウザでoutput.htmlを開いて確認してください。")

if __name__ == "__main__":
    main()